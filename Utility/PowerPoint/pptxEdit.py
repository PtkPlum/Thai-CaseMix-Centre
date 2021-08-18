import pandas as pd
from pptx import Presentation 
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE 
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR


prs = Presentation("ReportH11688Sana.pptx")

nSlide= len(prs.slides)
print("nSlide =", nSlide)

i = 0
Page = prs.slides[i]
print("\nPage {} = {}\n".format(i, Page))

#Txt = Page.shapes.text_frame
AllShapes = [i for i in Page.shapes]
print(AllShapes)

AllShapes[1].text = "ผลการศึกษาต้นทุนรายโรค ปีที่ 4"

prs.save("Result.pptx")