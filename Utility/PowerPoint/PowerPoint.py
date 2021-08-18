import pandas as pd
from pptx import Presentation 
from pptx.util import Inches
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE 
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR

class PowerPoint():
    def __init__(self):
        self.prs = Presentation("BlankSlide\BlankSlide.pptx")
        self.CurrentLayoutPage = None
        self.CurrentPage = None

    def NewPage(self):        
        self.CurrentLayoutPage = self.prs.slide_layouts[6]   # layout = 5
        self.CurrentPage = self.prs.slides.add_slide(self.CurrentLayoutPage) 
        #shape = self.prs.slides[0].shapes[0]
        #print(shape._element.xml)
        
    def Title(self, Text):
        self.TextBox(Text, FontSize=Pt(70), layout="center", left = Inches(16)/2, top=Inches(3))
        


    # LocLf   = location from left (inches)
    # LocTp   = location from top  (inches)
    def Table(self, DataFrame, LocLf=Inches(0.0), LocTp=Inches(1.0), width=Inches(6.0) , height=Inches(0.8)):
        FieldName = [i for i in DataFrame]
        nRow = 1 + len(DataFrame[FieldName[0]])
        nColum = len(FieldName) 
        width = height = 0
        table = self.CurrentPage.shapes.add_table(nRow, nColum, LocLf, LocTp, width, height).table
        # set column widths
        # write column headings
        MaxSizeList = []
        for i in range(nColum):
            SizeList = []
            field = ""
            for j in range(nRow - 1):
                table.cell(j+1, i).text = str( DataFrame[FieldName[i]][j] )
                SizeList.append(len(str( DataFrame[FieldName[i]][j] )))
            MaxSizeList.append(max(SizeList))
            if len(FieldName[i]) > MaxSizeList[-1]:
                count = 0
                for k in FieldName[i]:
                    if count+1 > MaxSizeList[-1] and k.isupper():
                        field += "\n" + k
                        count = 0
                    else:
                        field += k
                        count += 1
            else:
                field = FieldName[i] 
            table.cell(0, i).text = field

        
        for cell in self.iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "tahoma"
                    run.font.size = Pt(19)
                    
        ChgWidColList = [i for i in range(nColum)]
        for i in ChgWidColList:
            table.columns[i].width = Inches(1.6)
        '''
        ChgWidColList = [i for i in range(nColum)]
        CopyList = ChgWidColList.copy()
        UsedWidth = 0
        for i in ChgWidColList:
            if MaxSizeList[i] <= 5:
                table.columns[i].width = Inches(0.9)
                UsedWidth += Inches(0.9)
            elif MaxSizeList[i] <= 10 and MaxSizeList[i] > 4:
                table.columns[i].width = Inches(1.4)
                UsedWidth += Inches(1.4)
            elif MaxSizeList[i] <= 15 and MaxSizeList[i] > 10:
                table.columns[i].width = Inches(1.6)
                UsedWidth += Inches(1.6)
            else:
                CopyList.remove(i)
            
        ChgWidColList = CopyList.copy()
                
        for i in ChgWidColList:
            table.columns[i].width = Inches((16 - (UsedWidth + LocLf) / 914400) / len(ChgWidColList)) 
        '''    
        
        for j in range(nRow):
            table.rows[j].height = Inches((9 - (LocTp / 914400)) / nRow)
        
        
        table.cell(0, 0).merge(table.cell(0, 1))

    def iter_cells(self, table):
        for row in table.rows:
            for cell in row.cells:
                yield cell
                
    def TextBox(self, text, FontName="Calibri", FontSize=Pt(40), layout="left", 
                left=Inches(3.0), top=0, height=Inches(3.0), color=[0,0,0]):
        #width=Inches(len(text)*0.25)
        height = 0
        width = 0
        #left = (Inches(16) - width)/2
        txBox = self.CurrentPage.shapes.add_textbox(left, top, width, height) 
        tf = txBox.text_frame
        tf.text = ""
        p = tf.add_paragraph()
    
        run = p.add_run()
        run.text = text
        p.font.name = FontName
        p.font.size = FontSize
        run.font.color.rgb = RGBColor(color[0], color[1], color[2])
        if layout == "center":
            p.alignment = PP_ALIGN.CENTER
        elif layout == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            if layout != "left":
                print("ระบบไม่มี layout ", layout, ", กำหนด layout เป็น left แทน")
            p.alignment = PP_ALIGN.LEFT
        
        run = p.add_run()
        #run.text = " This is color you inputted"
        run.font.color.rgb = RGBColor(color[0], color[1], color[2])

            
    def Image(self, Pic, Layer = -1, left=Inches(2.0), right = Inches(3), width=Inches(7)):
        add_picture = self.CurrentPage.shapes.add_picture(Pic, left, right, width)
        
        self.CurrentPage.shapes._spTree.remove(add_picture._element)
        self.CurrentPage.shapes._spTree.insert(2, add_picture._element)


    def AutoShape(self, Type="ACTION_BUTTON_CUSTOM", Col="blueTCMC", brightness=0, left=Inches(3.0), top=Inches(1.0), 
                  width=Inches(2.0), height=Inches(3.0), text="", FontName="Calibri", FontSize=Pt(40), layout="left", 
                  color=[0,0,0]): #
        shapes = self.CurrentPage.shapes        
        shape = shapes.add_shape(MSO_SHAPE.ACTION_BUTTON_CUSTOM, left, top, width, height) 
        fill = shape.fill
        fill.solid()
        if Col=="blueTCMC":
            fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
        elif Col=="redPantone":
            fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2
        elif Col=="greenTea":
            fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_3
        elif Col=="purple":
            fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_4
        elif Col=="blueSea":
            fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_5
        elif Col=="orange":
            fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_6
        elif Col == "black":
            fill.fore_color.theme_color = MSO_THEME_COLOR.DARK_1
        else:
            fill.fore_color.theme_color = MSO_THEME_COLOR.DARK_1
            print("ไม่มี Col ที่พิมพ์มา ระบบเปลี่ยนเป็น black")
        fill.fore_color.brightness = brightness
        self.TextBox(text, FontName=FontName, FontSize=FontSize, layout=layout, left=left, top=top, height=height, color=color)
        


    def Save(self, name='Result'):
        Name = "Result\\"
        Name += name + ".pptx"
        self.prs.save(Name)
        print("\nไฟล์", name + ".pptx ได้ถูกเซฟเรียบร้อยแล้วครับ")